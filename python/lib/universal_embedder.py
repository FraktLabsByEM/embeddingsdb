import io
import os
import base64
import decord
import tempfile
import soundfile as sf
import torch
import whisper
import librosa
import numpy as np
import pytesseract
import pandas as pd
from PIL import Image
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
from pdf2image import convert_from_bytes
from moviepy.video.io.VideoFileClip import VideoFileClip
from imagebind.models import imagebind_model
from imagebind.models.imagebind_model import ModalityType
from imagebind.data import load_and_transform_audio_data, load_and_transform_vision_data, load_and_transform_text

print("Imports ok!")

plain_text_types = [
        "text/plain",                      # TXT
        "application/json",                 # JSON
        "application/xml", "text/xml",      # XML
        "text/html",                        # HTML
        "text/css",                         # CSS
        "text/javascript",                   # JavaScript
        "application/javascript",            # JavaScript (alternativo)
        "application/x-httpd-php",           # PHP
        "application/x-sh",                  # Shell script
        "application/x-python-code",         # Python
        "text/markdown",                     # Markdown
        "application/sql",                   # SQL scripts
        "text/csv",                          # CSV
        "application/x-yaml", "text/yaml",   # YAML
        "text/x-c", "text/x-c++",            # C y C++
        "text/x-java-source",                # Java
        "text/x-go",                         # Go
        "text/x-ruby",                       # Ruby
        "text/x-perl",                       # Perl
        "text/x-php",                        # PHP
        "text/x-shellscript"                 # Shell scripts (.sh)
    ]

image_types = [
        "image/jpeg",
        "image/png",
        "image/gif",
        "image/bmp",
        "image/tiff"
    ]

class UniversalEmbedder:
    def __init__(self):
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        # Init transcription model
        self.tr_model = whisper.load_model("small", device=self.device)
        # Configure tesserract
        self.tesseract_path = "/usr/bin/tesseract"  # Path to tesseract
        pytesseract.pytesseract.tesseract_cmd = self.tesseract_path
        # Imagebind model (512D shared embedding space for text, image, audio)
        print(f"Attempting to load embedding models into {self.device}")
        self.model = imagebind_model.imagebind_huge(pretrained=True).to(self.device)
        self.model.eval()
        
        self.audio_sample_rate = 44100

    def embed(self, data, storable = False):
        """
            Arguments:
                data(json request params):
                    - input: plain text or base64
            Returns:
                object(images do not return raw):
                    - emb: [0,1,2,3],
                    - raw: ["paragraph1", "paragraph1"]
        """
        
        try:
            # Text without base64 header is assumed to be plain input
            if not data["input"].startswith("data:") and ";base64," not in data["input"]:
                return self.embed_text(data["input"], storable)
            else:
                # Validate correct headers
                if "," not in data["input"]:
                    raise ValueError("embed - Invalid base64 input format. Base64 should include the header.")
                
                # Retrieve base64 header
                header, content = data["input"].split(",", 1)
                header = header.split(":")[1].split(";")[0]
                # Retrieve base64 content
                decoded_data = base64.b64decode(content)
                
                # Plain text types
                if header in plain_text_types:
                    print("Plain text file (TXT, JSON, XML, HTML, JS, CSS, etc.)")
                    text = decoded_data.decode("utf-8")  # Get content
                    return self.embed_text(text, storable)
                    
                # PDF document types
                elif "application/pdf" in header:
                    print("PDF file")
                    chunks = []
                    with io.BytesIO(decoded_data) as pdf_buffer:
                        reader = PdfReader(pdf_buffer)
                        for page in reader.pages:
                            page_text = page.extract_text()
                            if page_text:
                                page_text = page_text.replace("\n", "")
                                chunks.append(page_text)
                    
                    # Convert PDF pages to images for OCR
                    if len(chunks) == 0:
                        print("Applying OCR to scanned PDF")
                        images = convert_from_bytes(decoded_data, dpi=300)
                        for image in images:
                            ocr_text = self.ocr_image(image).strip()
                            chunks.append(ocr_text)
                    
                    # Build results
                    results = { "bytes": [], "raw": [] }
                    for chunk in chunks:
                        ch_res = self.embed_text(chunk, storable)
                        for r, b in zip(ch_res["raw"], ch_res["bytes"]):
                            results["raw"].append(r)
                            results["bytes"].append(b)
                                
                    return results

                # Word (DOCX)
                elif "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in header:
                    print("Word File (DOCX)")
                    
                    # Read document
                    with io.BytesIO(decoded_data) as doc_buffer:
                        doc = Document(doc_buffer)  # Load .docx file
                        chunks = [p.text.strip() for p in doc.paragraphs if p.text.strip()]  # Extract non empty paragraphs
                    
                    # Build results
                    results = { "bytes": [], "raw": [] }
                    for chunk in chunks:
                        ch_res = self.embed_text(chunk, storable)
                        for r, b in zip(ch_res["raw"], ch_res["bytes"]):
                            results["raw"].append(r)
                            results["bytes"].append(b)
                    
                    return results

                # Excel (XLSX, XLS)
                elif header in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"]:
                    print("Excel file (XLSX, XLS)")
                    
                    with io.BytesIO(decoded_data) as file_buffer:
                        xls = pd.ExcelFile(file_buffer)  # Load the Excel file
                        chunks = [xls.parse(sheet).to_string(index=False) for sheet in xls.sheet_names]  # Convert each sheet to a string
                    
                    
                    # Build results
                    results = { "bytes": [], "raw": [] }
                    for chunk in chunks:
                        ch_res = self.embed_text(chunk, storable)
                        for r, b in zip(ch_res["raw"], ch_res["bytes"]):
                            results["raw"].append(r)
                            results["bytes"].append(b)
                                
                    return results

                # PowerPoint (PPTX, PPT)
                elif header in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
                    print("PowerPoint file (PPTX, PPT)")
                    
                    with io.BytesIO(decoded_data) as file_buffer:
                        presentation = Presentation(file_buffer)  # Load the PowerPoint file
                        chunks = []
                        
                        for slide in presentation.slides:
                            slide_text = "\n".join([shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()])
                            if slide_text:
                                chunks.append(slide_text)  # Each slide is a separate chunk
                    
                    
                    # Build results
                    results = { "bytes": [], "raw": [] }
                    for chunk in chunks:
                        ch_res = self.embed_text(chunk, storable)
                        for r, b in zip(ch_res["raw"], ch_res["bytes"]):
                            results["raw"].append(r)
                            results["bytes"].append(b)
                                
                    return results

                # Imágenes (JPG, PNG, GIF, BMP, TIFF)
                elif header in image_types:
                    print("Image file (JPG, PNG, GIF, BMP, TIFF)")
                    
                    result = { "bytes": [], "raw": [] }
                    
                    with io.BytesIO(decoded_data) as image_buffer:
                        # Attempt to extract image embeddings
                        image_embedding = self.embed_image(image_buffer)
                        result["raw"].append("Protected image data")
                        result["bytes"].append(image_embedding.tolist() if storable else image_embedding)
                        # Attempt to extract text using OCR
                        ocr_text = self.ocr_image(image_buffer).strip()
                        ocr_embedings = self.embed_text(ocr_text, storable)
                        for r, b in zip(ocr_embedings["raw"], ocr_embedings["bytes"]):
                            result["raw"].append(r)
                            result["bytes"].append(b)
                    
                    return result

                # Audios (MP3, WAV, OGG, AAC)
                elif header in ["audio/mpeg", "audio/wav", "audio/ogg", "audio/aac"]:
                    print("Audio file (MP3, WAV, OGG, AAC)")
                    
                    # Read file
                    with io.BytesIO(decoded_data) as audio_buffer:
                        audio, sr = librosa.load(audio_buffer, sr=self.audio_sample_rate)  # Load the audio file
                    
                    result = { "raw": [], "bytes": [] }
                    # Generate raw audio embeddings
                    audio_embedding = self.embed_audio(audio, sr)  # Generate audio embedding
                    result["raw"].append("Protected audio data")
                    result["bytes"].append(audio_embedding.tolist() if storable else audio_embedding)

                    # Transcribe speech to text using Whisper
                    with tempfile.NamedTemporaryFile(suffix=".wav", delete=True) as temp_audio:
                        sf.write(temp_audio.name, audio, sr)  # Save audio as WAV
                        transcribed_text = self.tr_model.transcribe(temp_audio.name)["text"].strip()
                        # Generate transcription embeddings
                        tr_embeddings = self.embed_text(transcribed_text, storable)
                        for r, b in zip(tr_embeddings["raw"], tr_embeddings["bytes"]):
                            result["raw"].append(r)
                            result["bytes"].append(b)
                    
                    return result
                    
                # Videos (MP4, AVI, MKV, WEBM, MOV)
                elif header in ["video/mp4", "video/x-msvideo", "video/x-matroska", "video/webm", "video/quicktime"]:
                    print("Video file (MP4, AVI, MKV, WEBM, MOV)")

                    try:
                        result = { "raw": [], "bytes": [] }
                        
                        # Attempt to read video file
                        with tempfile.NamedTemporaryFile(suffix=".mp4", delete=False) as temp_video:
                            temp_video.write(decoded_data)
                            temp_video.flush()  # Ensure data is written before reading
                            temp_video_path = temp_video.name  # Store file path
                            
                            
                            # Extract 1 frame each second
                            vr = decord.VideoReader(temp_video_path)
                            fps = vr.get_avg_fps() # Get framerate
                            frame_indexes = [int(i * fps) for i in range(int(len(vr) / fps))] # Get frame indexes
                            frames = [vr[i].asnumpy() for i in frame_indexes if i < len(vr)] # Get frames
                            
                            # Generate image embeddings per frame extracted
                            for frame in frames:
                                fr_emb = self.embed_image(frame)
                                result["raw"].append("Protected image data")
                                result["bytes"].append(fr_emb.tolist() if storable else fr_emb)
                            
                            # Extract audio embeddings
                            audio_path = temp_video_path.replace(".mp4", ".wav")
                            video = VideoFileClip(temp_video_path)
                            video.audio.write_audiofile(audio_path, codec="pcm_s16le", fps=self.audio_sample_rate)
                            
                            # Generate audio embeddings for extracted audio
                            audio_loaded, sr = librosa.load(audio_path, sr=self.audio_sample_rate)
                            audio_emb = self.embed_audio(audio_loaded, sr)
                            result["raw"].append("Protected audio data")
                            result["bytes"].append(audio_emb.tolist() if storable else audio_emb)
                            
                            # Extract text transcription for audio
                            transcribed_text = self.tr_model.transcribe(audio_path)["text"].strip()
                            
                            # Generate text embeddings for audio transcription
                            text_emb = self.embed_text(transcribed_text, storable)
                            for r, b in zip(text_emb["raw"], text_emb["bytes"]):
                                result["raw"].append(r)
                                result["bytes"].append(b)

                        return result
                    finally:
                        if os.path.exists(temp_video_path):
                            os.remove(temp_video_path)
                        if os.path.exists(audio_path):
                            os.remove(audio_path)

                # Unknown file type
                else:
                    raise ValueError("embed - Unknow file type.")
                    
        except Exception as e:
            raise ValueError(f"embed - Error processing file: {e}")


    def embed_text(self, text, storable):
        """Generate 512D text embeddings using AudioCLIP"""
        try:
            chunks = text.split("\n")
            response = { "bytes": [], "raw": [] }
            for chunk in chunks:
                chunk = chunk.strip()
                if chunk:
                    text_input = load_and_transform_text([chunk], self.device)
                    with torch.no_grad():
                        emb = self.model({ModalityType.TEXT: text_input})[ModalityType.TEXT]
                    emb = emb[0].cpu().numpy().flatten()
                    response["raw"].append(chunk)
                    response["bytes"].append(emb.tolist() if storable else emb)
            return response
        except Exception as err:
            raise ValueError(f"embed_text - Failed to generate text embeddings due to: {err}")


    def embed_image(self, image_bytes):
        try:
            if isinstance(image_bytes, bytes):
                img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
            elif isinstance(image_bytes, np.ndarray):
                img = Image.fromarray(image_bytes)
            else:
                raise ValueError("embed_image - Unsupported image format")

            with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp:
                img.save(tmp.name)
                image_input = load_and_transform_vision_data([tmp.name], self.device)

            with torch.no_grad():
                emb = self.model({ModalityType.VISION: image_input})[ModalityType.VISION]

            return emb[0].cpu().numpy().flatten()

        except Exception as err:
            raise ValueError(f"embed_image - Failed using ImageBind: {err}")


    def embed_audio(self, audio_np_array, sr):
        try:
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
                sf.write(tmp.name, audio_np_array, sr)
                audio_input = load_and_transform_audio_data([tmp.name], self.device)

            with torch.no_grad():
                emb = self.model({ModalityType.AUDIO: audio_input})[ModalityType.AUDIO]

            return emb[0].cpu().numpy().flatten()

        except Exception as err:
            raise ValueError(f"embed_audio - Failed using ImageBind: {err}")


    def ocr_image(self, image_bytes):
        """Apply OCR to provided image"""
        try:
            img = Image.open(image_bytes)
            return pytesseract.image_to_string(img)
        except Exception as err:
            raise ValueError(f"ocr - Failed to OCR image due to: {err}")
            
